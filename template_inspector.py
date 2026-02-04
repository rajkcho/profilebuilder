"""
Helper: iterate through every slide layout in a .pptx and print
each shape's placeholder idx so you can map data correctly.

Usage:  python template_inspector.py [path_to_pptx]
"""

import sys
from pptx import Presentation


def inspect(path: str = "assets/template.pptx"):
    prs = Presentation(path)
    print(f"Slide width : {prs.slide_width}")
    print(f"Slide height: {prs.slide_height}")
    print(f"Layouts     : {len(prs.slide_layouts)}\n")

    for i, layout in enumerate(prs.slide_layouts):
        print(f"── Layout {i}: {layout.name} ──")
        for shape in layout.placeholders:
            fmt = shape.placeholder_format
            print(
                f"   idx={fmt.idx:<4}  type={fmt.type!s:<20}  "
                f"name={shape.name!r:<25}  "
                f"pos=({shape.left}, {shape.top})  "
                f"size=({shape.width}x{shape.height})"
            )
        if not layout.placeholders:
            print("   (no placeholders)")
        print()


if __name__ == "__main__":
    path = sys.argv[1] if len(sys.argv) > 1 else "assets/template.pptx"
    inspect(path)
