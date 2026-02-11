"""
PPTX Generator — Goldman Sachs-style investment banking presentations.

Design: Clean, minimal, data-dense. Navy/white palette with gold accents.
Arial font throughout, tight grid alignment, professional formatting.
3-slide decks for both Company Profile and Merger Analysis.
"""

import io
from datetime import datetime
from typing import Optional

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.dates as mdates
import matplotlib.ticker as mticker
import numpy as np
import pandas as pd

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

from data_engine import (CompanyData, format_number, format_pct, format_multiple,
                         calculate_piotroski_score, calculate_intrinsic_value, get_key_ratios_summary)

# ══════════════════════════════════════════════════════════════
# GOLDMAN SACHS STYLE PALETTE
# ══════════════════════════════════════════════════════════════

# Modern fintech palette (Linear/Vercel/Stripe-inspired)
NAVY = RGBColor(0x0C, 0x0F, 0x1A)  # Deep charcoal background
LIGHT_NAVY = RGBColor(0x1F, 0x29, 0x37)  # Lighter dark
BLUE = RGBColor(0x25, 0x63, 0xEB)  # Electric blue accent (#2563EB)
EMERALD = RGBColor(0x10, 0xB9, 0x81)  # Emerald accent (#10B981)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)
LIGHT_GRAY = RGBColor(0x11, 0x18, 0x27)  # Dark card bg
MED_GRAY = RGBColor(0x37, 0x41, 0x55)
DARK_GRAY = RGBColor(0xE5, 0xE7, 0xEB)  # Light text on dark
TEXT_DIM = RGBColor(0x9C, 0xA3, 0xAF)
GREEN = RGBColor(0x10, 0xB9, 0x81)
RED = RGBColor(0xEF, 0x44, 0x44)
# Legacy aliases for compatibility
GOLD = BLUE
PURPLE = BLUE
PURPLE_LIGHT = RGBColor(0x60, 0xA5, 0xFA)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


# ══════════════════════════════════════════════════════════════
# CORE HELPERS
# ══════════════════════════════════════════════════════════════

def _set_cell_text(cell, text, font_size=9, bold=False, color=DARK_GRAY, align=PP_ALIGN.LEFT):
    """Set cell text with consistent formatting."""
    cell.text = str(text) if text is not None else "—"
    for p in cell.text_frame.paragraphs:
        p.font.size = Pt(font_size)
        p.font.bold = bold
        p.font.name = "Arial"
        p.font.color.rgb = color
        p.alignment = align
    cell.text_frame.paragraphs[0].space_before = Pt(2)
    cell.text_frame.paragraphs[0].space_after = Pt(2)


def _add_textbox(slide, left, top, width, height, text,
                 font_size=10, bold=False, color=DARK_GRAY, align=PP_ALIGN.LEFT):
    """Add a textbox with specified formatting."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = str(text) if text else ""
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.name = "Arial"
    p.font.color.rgb = color
    p.alignment = align
    return txBox


def _add_rect(slide, left, top, width, height, fill_color, line_color=None):
    """Add a rectangle shape."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape


def _gs_header(slide, title, subtitle=""):
    """Dark theme header with purple accent line."""
    # Dark background fill for whole slide
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = NAVY
    # Purple accent line at very top
    _add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.06), PURPLE)
    # Title
    _add_textbox(slide, Inches(0.5), Inches(0.25), Inches(10), Inches(0.4),
                 title.upper(), font_size=14, bold=True, color=WHITE)
    if subtitle:
        _add_textbox(slide, Inches(0.5), Inches(0.55), Inches(10), Inches(0.25),
                     subtitle, font_size=9, color=TEXT_DIM)
    # Thin line under title
    _add_rect(slide, Inches(0.5), Inches(0.85), Inches(12.333), Inches(0.015), PURPLE)


def _gs_footer(slide, left_text, right_text="ORBITAL", slide_num=None):
    """Dark theme footer with optional slide number."""
    # Thin line above footer
    _add_rect(slide, Inches(0.5), SLIDE_H - Inches(0.5), Inches(12.333), Inches(0.01), MED_GRAY)
    # Left text (confidential + date)
    _add_textbox(slide, Inches(0.5), SLIDE_H - Inches(0.4), Inches(8), Inches(0.3),
                 f"CONFIDENTIAL  |  {left_text}  |  {datetime.now().strftime('%B %Y')}",
                 font_size=7, color=TEXT_DIM)
    # Slide number (center)
    if slide_num is not None:
        _add_textbox(slide, Inches(6), SLIDE_H - Inches(0.4), Inches(1.333), Inches(0.3),
                     str(slide_num), font_size=7, color=TEXT_DIM, align=PP_ALIGN.CENTER)
    # Right text (brand)
    _add_textbox(slide, Inches(10), SLIDE_H - Inches(0.4), Inches(2.833), Inches(0.3),
                 right_text, font_size=7, bold=True, color=PURPLE, align=PP_ALIGN.RIGHT)


def _gs_section_title(slide, text, top):
    """Add a section title with purple underline."""
    _add_textbox(slide, Inches(0.5), top, Inches(5), Inches(0.3),
                 text.upper(), font_size=9, bold=True, color=PURPLE_LIGHT)
    _add_rect(slide, Inches(0.5), top + Inches(0.25), Inches(1.5), Inches(0.02), PURPLE)


def _gs_table(slide, headers, rows, left, top, width, height, col_widths=None):
    """Professional Goldman-style table."""
    n_rows = len(rows) + 1
    n_cols = len(headers)
    shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
    table = shape.table

    if col_widths:
        for i, w in enumerate(col_widths):
            if i < n_cols:
                table.columns[i].width = w

    # Header row - purple background
    for c, hdr in enumerate(headers):
        cell = table.cell(0, c)
        _set_cell_text(cell, hdr, font_size=8, bold=True, color=WHITE,
                       align=PP_ALIGN.CENTER if c > 0 else PP_ALIGN.LEFT)
        cell.fill.solid()
        cell.fill.fore_color.rgb = PURPLE

    # Data rows - dark theme alternating
    for r, row_data in enumerate(rows, start=1):
        for c, val in enumerate(row_data):
            cell = table.cell(r, c)
            is_first_col = (c == 0)
            _set_cell_text(cell, val, font_size=8, bold=is_first_col,
                           color=DARK_GRAY, align=PP_ALIGN.LEFT if is_first_col else PP_ALIGN.RIGHT)
            cell.fill.solid()
            cell.fill.fore_color.rgb = LIGHT_GRAY if r % 2 == 0 else RGBColor(0x12, 0x15, 0x28)

    return shape


def _year_labels(series, count=4):
    """Extract year labels from a pandas Series."""
    if series is None or len(series) == 0:
        return ["—"] * count
    labels = []
    for idx in series.index[:count]:
        if hasattr(idx, "year"):
            labels.append(str(idx.year))
        else:
            labels.append(str(idx))
    while len(labels) < count:
        labels.append("—")
    return labels


def _series_vals(series, count=4, currency_symbol="$"):
    """Format series values for display."""
    if series is None or len(series) == 0:
        return ["—"] * count
    vals = []
    for v in list(series.values[:count]):
        if v is not None and not pd.isna(v):
            vals.append(format_number(v, currency_symbol=currency_symbol))
        else:
            vals.append("—")
    while len(vals) < count:
        vals.append("—")
    return vals


# ══════════════════════════════════════════════════════════════
# COMPANY PROFILE - 3 SLIDES
# ══════════════════════════════════════════════════════════════

def _company_slide_1(prs, cd: CompanyData):
    """Slide 1: Executive Summary - Company overview + key metrics."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    cs = cd.currency_symbol

    _gs_header(slide, f"{cd.name} ({cd.ticker})", f"{cd.sector} | {cd.industry}")
    _gs_footer(slide, cd.ticker)

    # Left column: Company Description
    _gs_section_title(slide, "Company Overview", Inches(1.0))

    _raw_desc = getattr(cd, 'description', None) or getattr(cd, 'long_business_summary', None) or ""
    if hasattr(_raw_desc, 'iloc'):
        _raw_desc = str(_raw_desc.iloc[0]) if len(_raw_desc) > 0 else ""
    _raw_desc = str(_raw_desc) if _raw_desc else ""
    desc = (_raw_desc[:600] + "...") if len(_raw_desc) > 600 else (_raw_desc or "Company description not available.")
    _add_textbox(slide, Inches(0.5), Inches(1.4), Inches(5.8), Inches(2.0),
                 desc, font_size=9, color=DARK_GRAY)

    # Key Statistics box
    _gs_section_title(slide, "Key Statistics", Inches(3.6))
    stats = [
        ["Market Cap", format_number(cd.market_cap, currency_symbol=cs)],
        ["Enterprise Value", format_number(cd.enterprise_value, currency_symbol=cs)],
        ["Revenue (LTM)", format_number(cd.revenue.iloc[0] if cd.revenue is not None and len(cd.revenue) > 0 else None, currency_symbol=cs)],
        ["EBITDA (LTM)", format_number(cd.ebitda.iloc[0] if cd.ebitda is not None and len(cd.ebitda) > 0 else None, currency_symbol=cs)],
        ["Employees", f"{cd.employees:,}" if cd.employees else "—"],
    ]
    _gs_table(slide, ["Metric", "Value"], stats,
              Inches(0.5), Inches(3.95), Inches(5.8), Inches(1.8),
              col_widths=[Inches(2.5), Inches(3.3)])

    # Right column: Valuation & Trading
    _gs_section_title(slide, "Valuation Metrics", Inches(1.0))

    valuation = [
        ["Current Price", f"{cs}{cd.current_price:,.2f}" if cd.current_price else "—"],
        ["52-Week Range", f"{cs}{cd.fifty_two_week_low:,.2f} - {cs}{cd.fifty_two_week_high:,.2f}" if cd.fifty_two_week_low and cd.fifty_two_week_high else "—"],
        ["P/E (TTM)", f"{cd.trailing_pe:.1f}x" if cd.trailing_pe else "—"],
        ["P/E (Forward)", f"{cd.forward_pe:.1f}x" if cd.forward_pe else "—"],
        ["EV/EBITDA", f"{cd.ev_to_ebitda:.1f}x" if cd.ev_to_ebitda else "—"],
        ["EV/Revenue", f"{cd.ev_to_revenue:.1f}x" if cd.ev_to_revenue else "—"],
        ["P/B", f"{cd.price_to_book:.1f}x" if cd.price_to_book else "—"],
        ["Dividend Yield", f"{cd.dividend_yield*100:.2f}%" if cd.dividend_yield else "—"],
    ]
    _gs_table(slide, ["Metric", "Value"], valuation,
              Inches(6.8), Inches(1.35), Inches(6), Inches(2.8),
              col_widths=[Inches(3), Inches(3)])

    # Trading Info
    _gs_section_title(slide, "Trading Information", Inches(4.4))
    trading = [
        ["Exchange", cd.exchange or "—"],
        ["Beta", f"{cd.beta:.2f}" if cd.beta else "—"],
        ["Avg Volume (3M)", f"{cd.average_volume/1e6:.1f}M" if cd.average_volume else "—"],
    ]
    _gs_table(slide, ["Metric", "Value"], trading,
              Inches(6.8), Inches(4.75), Inches(6), Inches(1.1),
              col_widths=[Inches(3), Inches(3)])


def _company_slide_2(prs, cd: CompanyData):
    """Slide 2: Financial Overview - P&L, Balance Sheet, Cash Flow."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    cs = cd.currency_symbol

    _gs_header(slide, "Financial Overview", f"{cd.name} ({cd.ticker})")
    _gs_footer(slide, cd.ticker)

    years = _year_labels(cd.revenue, 4)

    # Income Statement (top left)
    _gs_section_title(slide, "Income Statement", Inches(1.0))
    income_rows = [
        ["Revenue"] + _series_vals(cd.revenue, 4, cs),
        ["Gross Profit"] + _series_vals(cd.gross_profit, 4, cs),
        ["EBITDA"] + _series_vals(cd.ebitda, 4, cs),
        ["Net Income"] + _series_vals(cd.net_income, 4, cs),
    ]
    _gs_table(slide, [""] + years, income_rows,
              Inches(0.5), Inches(1.35), Inches(6), Inches(1.5),
              col_widths=[Inches(1.5)] + [Inches(1.125)] * 4)

    # Margins (top right)
    _gs_section_title(slide, "Profitability Margins", Inches(1.0))

    def _margin_vals(series, count=4):
        if series is None or len(series) == 0:
            return ["—"] * count
        vals = []
        for v in list(series.values[:count]):
            if v is not None and not pd.isna(v):
                vals.append(f"{v*100:.1f}%")
            else:
                vals.append("—")
        while len(vals) < count:
            vals.append("—")
        return vals

    margin_rows = [
        ["Gross Margin"] + _margin_vals(cd.gross_margin_series, 4),
        ["EBITDA Margin"] + _margin_vals(cd.ebitda_margin, 4),
        ["Net Margin"] + _margin_vals(cd.net_margin_series, 4),
    ]
    _gs_table(slide, [""] + years, margin_rows,
              Inches(6.8), Inches(1.35), Inches(6), Inches(1.2),
              col_widths=[Inches(1.5)] + [Inches(1.125)] * 4)

    # Balance Sheet (middle left)
    _gs_section_title(slide, "Balance Sheet", Inches(3.0))
    bs_years = _year_labels(cd.total_assets, 4)
    bs_rows = [
        ["Total Assets"] + _series_vals(cd.total_assets, 4, cs),
        ["Cash & Equivalents"] + _series_vals(cd.cash_and_equivalents, 4, cs),
        ["Total Debt"] + _series_vals(cd.total_debt, 4, cs),
        ["Total Equity"] + _series_vals(cd.total_equity, 4, cs),
    ]
    _gs_table(slide, [""] + bs_years, bs_rows,
              Inches(0.5), Inches(3.35), Inches(6), Inches(1.5),
              col_widths=[Inches(1.5)] + [Inches(1.125)] * 4)

    # Cash Flow (middle right)
    _gs_section_title(slide, "Cash Flow Statement", Inches(3.0))
    cf_years = _year_labels(cd.operating_cashflow_series, 4)
    cf_rows = [
        ["Operating CF"] + _series_vals(cd.operating_cashflow_series, 4, cs),
        ["Capital Expenditures"] + _series_vals(cd.capital_expenditure, 4, cs),
        ["Free Cash Flow"] + _series_vals(cd.free_cashflow_series, 4, cs),
    ]
    _gs_table(slide, [""] + cf_years, cf_rows,
              Inches(6.8), Inches(3.35), Inches(6), Inches(1.2),
              col_widths=[Inches(1.5)] + [Inches(1.125)] * 4)

    # Credit Metrics (bottom)
    _gs_section_title(slide, "Credit Metrics", Inches(5.0))

    debt_ebitda = "—"
    if cd.total_debt is not None and len(cd.total_debt) > 0 and cd.ebitda is not None and len(cd.ebitda) > 0:
        d = cd.total_debt.iloc[0]
        e = cd.ebitda.iloc[0]
        if d and e and e != 0:
            debt_ebitda = f"{d/e:.1f}x"

    credit_rows = [
        ["Debt / EBITDA", debt_ebitda],
        ["Net Debt / EBITDA", f"{cd.net_debt_to_ebitda:.1f}x" if cd.net_debt_to_ebitda else "—"],
        ["Interest Coverage", f"{cd.interest_coverage:.1f}x" if cd.interest_coverage else "—"],
    ]
    _gs_table(slide, ["Metric", "Value"], credit_rows,
              Inches(0.5), Inches(5.35), Inches(4), Inches(1.1),
              col_widths=[Inches(2), Inches(2)])

    # Piotroski F-Score (bottom right)
    try:
        piotroski = calculate_piotroski_score(cd)
        if piotroski:
            _gs_section_title(slide, "Piotroski F-Score", Inches(5.0))
            score = piotroski['score']
            score_color = GREEN if score >= 7 else (RED if score <= 3 else DARK_GRAY)
            label = "Strong" if score >= 7 else ("Weak" if score <= 3 else "Neutral")
            _add_textbox(slide, Inches(6.8), Inches(5.35), Inches(2), Inches(0.5),
                         f"{score} / 9  ({label})", font_size=16, bold=True, color=score_color)
    except Exception:
        pass


def _company_slide_3(prs, cd: CompanyData):
    """Slide 3: Peer Comparison & Analyst Views."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    cs = cd.currency_symbol

    _gs_header(slide, "Valuation & Peer Comparison", f"{cd.name} ({cd.ticker})")
    _gs_footer(slide, cd.ticker)

    # Peer Comparison Table
    _gs_section_title(slide, "Comparable Company Analysis", Inches(1.0))

    peer_headers = ["Company", "Price", "Mkt Cap", "EV/EBITDA", "P/E", "EV/Rev"]
    peer_rows = []

    # Add target company first
    peer_rows.append([
        f"{cd.ticker} (Target)",
        f"{cs}{cd.current_price:.2f}" if cd.current_price else "—",
        format_number(cd.market_cap, currency_symbol=cs) if cd.market_cap else "—",
        f"{cd.ev_to_ebitda:.1f}x" if cd.ev_to_ebitda else "—",
        f"{cd.trailing_pe:.1f}x" if cd.trailing_pe else "—",
        f"{cd.ev_to_revenue:.1f}x" if cd.ev_to_revenue else "—",
    ])

    # Add peers
    if cd.peer_data:
        for peer in cd.peer_data[:6]:
            peer_rows.append([
                peer.get("ticker", "—"),
                f"{cs}{peer.get('price', 0):.2f}" if peer.get("price") else "—",
                format_number(peer.get("market_cap"), currency_symbol=cs) if peer.get("market_cap") else "—",
                f"{peer.get('ev_ebitda', 0):.1f}x" if peer.get("ev_ebitda") else "—",
                f"{peer.get('pe_ratio', 0):.1f}x" if peer.get("pe_ratio") else "—",
                f"{peer.get('ev_revenue', 0):.1f}x" if peer.get("ev_revenue") else "—",
            ])

    _gs_table(slide, peer_headers, peer_rows,
              Inches(0.5), Inches(1.35), Inches(12.333), Inches(2.5),
              col_widths=[Inches(2.5), Inches(1.5), Inches(2.2), Inches(2), Inches(2), Inches(2.133)])

    # Analyst Price Targets (bottom left)
    _gs_section_title(slide, "Analyst Price Targets", Inches(4.2))

    if cd.analyst_price_targets:
        targets = cd.analyst_price_targets
        target_rows = [
            ["Current Price", f"{cs}{cd.current_price:.2f}" if cd.current_price else "—"],
            ["Low Target", f"{cs}{targets.get('low', 0):.2f}" if targets.get('low') else "—"],
            ["Mean Target", f"{cs}{targets.get('mean', 0):.2f}" if targets.get('mean') else "—"],
            ["High Target", f"{cs}{targets.get('high', 0):.2f}" if targets.get('high') else "—"],
        ]
        if targets.get('mean') and cd.current_price:
            upside = ((targets['mean'] / cd.current_price) - 1) * 100
            target_rows.append(["Implied Upside", f"{upside:+.1f}%"])
    else:
        target_rows = [["No analyst data available", "—"]]

    _gs_table(slide, ["Metric", "Value"], target_rows,
              Inches(0.5), Inches(4.55), Inches(5), Inches(1.8),
              col_widths=[Inches(2.5), Inches(2.5)])

    # Analyst Recommendations (bottom right)
    _gs_section_title(slide, "Analyst Recommendations", Inches(4.2))

    if cd.analyst_recommendations:
        rec = cd.analyst_recommendations
        rec_rows = [
            ["Strong Buy", str(rec.get('strongBuy', 0))],
            ["Buy", str(rec.get('buy', 0))],
            ["Hold", str(rec.get('hold', 0))],
            ["Sell", str(rec.get('sell', 0))],
            ["Strong Sell", str(rec.get('strongSell', 0))],
        ]
    else:
        rec_rows = [["No recommendation data available", "—"]]

    _gs_table(slide, ["Rating", "# Analysts"], rec_rows,
              Inches(6.8), Inches(4.55), Inches(5.533), Inches(1.8),
              col_widths=[Inches(2.8), Inches(2.733)])


def _company_slide_valuation(prs, cd: CompanyData):
    """Slide 4: Valuation Summary - Multiples, DCF, Investment Thesis."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    cs = cd.currency_symbol

    _gs_header(slide, "Valuation Summary", f"{cd.name} ({cd.ticker})")
    _gs_footer(slide, cd.ticker)

    # Key Multiples Table (top left)
    _gs_section_title(slide, "Key Valuation Multiples", Inches(1.0))

    multiples_rows = [
        ["P/E (TTM)", f"{cd.trailing_pe:.1f}x" if cd.trailing_pe else "—"],
        ["P/E (Forward)", f"{cd.forward_pe:.1f}x" if cd.forward_pe else "—"],
        ["EV/EBITDA", f"{cd.ev_to_ebitda:.1f}x" if cd.ev_to_ebitda else "—"],
        ["P/S", f"{cd.price_to_sales:.1f}x" if cd.price_to_sales else "—"],
        ["P/B", f"{cd.price_to_book:.1f}x" if cd.price_to_book else "—"],
        ["EV/Revenue", f"{cd.ev_to_revenue:.1f}x" if cd.ev_to_revenue else "—"],
        ["PEG Ratio", f"{cd.peg_ratio:.2f}" if cd.peg_ratio else "—"],
    ]

    # Add peer median column if peers available
    if cd.peer_data:
        headers = ["Multiple", cd.ticker, "Peer Median"]
        pe_vals = [p.get('trailing_pe') for p in cd.peer_data if p.get('trailing_pe')]
        fpe_vals = [p.get('forward_pe') for p in cd.peer_data if p.get('forward_pe')]
        ev_ebitda_vals = [p.get('ev_to_ebitda') for p in cd.peer_data if p.get('ev_to_ebitda')]
        ps_vals = [p.get('price_to_sales') for p in cd.peer_data if p.get('price_to_sales')]
        pb_vals = [p.get('price_to_book') for p in cd.peer_data if p.get('price_to_book')]
        peg_vals = [p.get('peg_ratio') for p in cd.peer_data if p.get('peg_ratio')]

        def _med(vals, fmt=":.1f"):
            if not vals:
                return "—"
            m = float(np.median(vals))
            return f"{m:.1f}x"

        multiples_rows = [
            ["P/E (TTM)", f"{cd.trailing_pe:.1f}x" if cd.trailing_pe else "—", _med(pe_vals)],
            ["P/E (Forward)", f"{cd.forward_pe:.1f}x" if cd.forward_pe else "—", _med(fpe_vals)],
            ["EV/EBITDA", f"{cd.ev_to_ebitda:.1f}x" if cd.ev_to_ebitda else "—", _med(ev_ebitda_vals)],
            ["P/S", f"{cd.price_to_sales:.1f}x" if cd.price_to_sales else "—", _med(ps_vals)],
            ["P/B", f"{cd.price_to_book:.1f}x" if cd.price_to_book else "—", _med(pb_vals)],
            ["EV/Revenue", f"{cd.ev_to_revenue:.1f}x" if cd.ev_to_revenue else "—", "—"],
            ["PEG Ratio", f"{cd.peg_ratio:.2f}" if cd.peg_ratio else "—", _med(peg_vals)],
        ]
        _gs_table(slide, headers, multiples_rows,
                  Inches(0.5), Inches(1.35), Inches(6), Inches(2.5),
                  col_widths=[Inches(2), Inches(2), Inches(2)])
    else:
        _gs_table(slide, ["Multiple", "Value"], multiples_rows,
                  Inches(0.5), Inches(1.35), Inches(5), Inches(2.5),
                  col_widths=[Inches(2.5), Inches(2.5)])

    # DCF / Intrinsic Value (top right)
    _gs_section_title(slide, "Intrinsic Value Estimate (DCF)", Inches(1.0))
    try:
        iv = calculate_intrinsic_value(cd)
        if iv:
            iv_rows = [
                ["Current Price", f"{cs}{cd.current_price:.2f}" if cd.current_price else "—"],
                ["Intrinsic Value / Share", f"{cs}{iv['intrinsic_value_per_share']:.2f}"],
                ["Upside / (Downside)", f"{iv['upside_pct']:+.1f}%"],
                ["Margin of Safety", f"{iv['margin_of_safety']:.1f}%"],
            ]
            _gs_table(slide, ["Metric", "Value"], iv_rows,
                      Inches(7), Inches(1.35), Inches(5.833), Inches(1.5),
                      col_widths=[Inches(3), Inches(2.833)])

            # Assumptions footnote
            a = iv['assumptions']
            _add_textbox(slide, Inches(7), Inches(2.9), Inches(5.8), Inches(0.4),
                         f"Assumptions: Growth {a['growth_rate']*100:.0f}%, Discount {a['discount_rate']*100:.0f}%, "
                         f"Terminal Multiple {a['terminal_multiple']}x, Base FCF {format_number(a['base_fcf'], currency_symbol=cs)}",
                         font_size=7, color=DARK_GRAY)
        else:
            _add_textbox(slide, Inches(7), Inches(1.35), Inches(5.8), Inches(0.5),
                         "Insufficient data for DCF analysis (negative or missing FCF)",
                         font_size=9, color=DARK_GRAY)
    except Exception:
        _add_textbox(slide, Inches(7), Inches(1.35), Inches(5.8), Inches(0.5),
                     "DCF analysis unavailable", font_size=9, color=DARK_GRAY)

    # Implied Price Range (middle)
    _gs_section_title(slide, "Implied Price Range", Inches(4.0))

    price_points = []
    if cd.analyst_price_targets:
        low_t = cd.analyst_price_targets.get('low')
        high_t = cd.analyst_price_targets.get('high')
        mean_t = cd.analyst_price_targets.get('mean')
        if low_t and high_t:
            price_points.append(["Analyst Range", f"{cs}{low_t:.2f} — {cs}{high_t:.2f}"])
        if mean_t:
            price_points.append(["Analyst Mean", f"{cs}{mean_t:.2f}"])
    try:
        iv = calculate_intrinsic_value(cd)
        if iv:
            price_points.append(["DCF Implied", f"{cs}{iv['intrinsic_value_per_share']:.2f}"])
    except Exception:
        pass
    if cd.current_price and cd.trailing_pe:
        price_points.append(["Current (Market)", f"{cs}{cd.current_price:.2f}"])

    if price_points:
        _gs_table(slide, ["Method", "Price"], price_points,
                  Inches(0.5), Inches(4.35), Inches(5), Inches(1.5),
                  col_widths=[Inches(2.5), Inches(2.5)])
    else:
        _add_textbox(slide, Inches(0.5), Inches(4.35), Inches(5), Inches(0.5),
                     "Insufficient data for price range analysis", font_size=9, color=DARK_GRAY)

    # Investment Thesis (bottom right)
    _gs_section_title(slide, "Investment Thesis", Inches(4.0))

    bullets = []
    # Auto-generate thesis bullets from data
    if cd.revenue_growth is not None:
        direction = "growing" if cd.revenue_growth > 0 else "declining"
        bullets.append(f"Revenue {direction} at {abs(cd.revenue_growth):.1f}% YoY")
    if cd.profit_margins is not None:
        bullets.append(f"Net margin of {cd.profit_margins*100:.1f}% {'above' if cd.profit_margins > 0.10 else 'below'} average")
    if cd.return_on_equity is not None and cd.return_on_equity > 0.15:
        bullets.append(f"Strong ROE of {cd.return_on_equity*100:.1f}% indicates competitive advantage")
    if cd.debt_to_equity is not None:
        lev = "low" if cd.debt_to_equity < 50 else ("moderate" if cd.debt_to_equity < 150 else "high")
        bullets.append(f"Leverage is {lev} with D/E of {cd.debt_to_equity:.0f}%")
    try:
        piotroski = calculate_piotroski_score(cd)
        if piotroski:
            s = piotroski['score']
            qual = "strong" if s >= 7 else ("weak" if s <= 3 else "moderate")
            bullets.append(f"Piotroski F-Score of {s}/9 signals {qual} financial health")
    except Exception:
        pass

    if not bullets:
        bullets = ["Insufficient data to generate investment thesis"]

    # Add bullets as textbox
    thesis_text = "\n".join(f"• {b}" for b in bullets[:4])
    _add_textbox(slide, Inches(7), Inches(4.35), Inches(5.833), Inches(2.5),
                 thesis_text, font_size=9, color=DARK_GRAY)


def _add_confidential_watermark(slide):
    """Add a diagonal 'CONFIDENTIAL' watermark to the slide."""
    from pptx.oxml.ns import qn
    txBox = slide.shapes.add_textbox(Inches(2), Inches(2.5), Inches(9), Inches(2.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "CONFIDENTIAL"
    p.font.size = Pt(60)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0x6B, 0x5C, 0xE7)
    p.alignment = PP_ALIGN.CENTER
    # Set rotation and transparency via XML
    sp = txBox._element
    sp.attrib['rot'] = '-2700000'  # -45 degrees in EMUs
    # Set transparency via solidFill alpha
    for run_elem in sp.iter(qn('a:solidFill')):
        srgb = run_elem.find(qn('a:srgbClr'))
        if srgb is not None:
            alpha = srgb.makeelement(qn('a:alpha'), {'val': '15000'})
            srgb.append(alpha)


def _add_slide_numbers(prs):
    """Add slide numbers to all slides."""
    for i, slide in enumerate(prs.slides, start=1):
        _add_textbox(slide, Inches(6), SLIDE_H - Inches(0.4), Inches(1.333), Inches(0.3),
                     str(i), font_size=7, color=TEXT_DIM, align=PP_ALIGN.CENTER)


def _company_slide_esg(prs, cd: CompanyData, slide_num=None):
    """ESG Summary slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    cs = cd.currency_symbol

    _gs_header(slide, "ESG Summary", f"{cd.name} ({cd.ticker})")
    _gs_footer(slide, cd.ticker, slide_num=slide_num)

    # ESG scores (use available data or show framework)
    _gs_section_title(slide, "Environmental, Social & Governance", Inches(1.0))

    esg_score = getattr(cd, 'esg_score', None)
    env_score = getattr(cd, 'environmental_score', None)
    soc_score = getattr(cd, 'social_score', None)
    gov_score = getattr(cd, 'governance_score', None)

    esg_rows = [
        ["Overall ESG Score", f"{esg_score:.0f}/100" if esg_score else "Not Rated"],
        ["Environmental", f"{env_score:.0f}/100" if env_score else "N/A"],
        ["Social", f"{soc_score:.0f}/100" if soc_score else "N/A"],
        ["Governance", f"{gov_score:.0f}/100" if gov_score else "N/A"],
    ]
    _gs_table(slide, ["Category", "Score"], esg_rows,
              Inches(0.5), Inches(1.35), Inches(5.5), Inches(1.5),
              col_widths=[Inches(2.8), Inches(2.7)])

    # ESG Risk Assessment
    _gs_section_title(slide, "ESG Risk Assessment", Inches(3.2))

    controversy_level = getattr(cd, 'controversy_level', None)
    risk_items = [
        ["Controversy Level", f"Level {controversy_level}" if controversy_level else "N/A"],
        ["Sector", cd.sector or "N/A"],
        ["Industry ESG Profile", "High Scrutiny" if cd.sector in ["Energy", "Utilities", "Basic Materials"] else "Standard"],
    ]
    _gs_table(slide, ["Risk Factor", "Assessment"], risk_items,
              Inches(0.5), Inches(3.55), Inches(5.5), Inches(1.1),
              col_widths=[Inches(2.8), Inches(2.7)])

    # ESG Framework (right side)
    _gs_section_title(slide, "ESG Integration Framework", Inches(1.0))
    framework_text = (
        "• Environmental: Carbon footprint, resource efficiency, waste management\n"
        "• Social: Employee relations, diversity, supply chain labor standards\n"
        "• Governance: Board independence, executive compensation, shareholder rights\n"
        "• Regulatory: Compliance with local and international ESG standards"
    )
    _add_textbox(slide, Inches(6.8), Inches(1.35), Inches(6), Inches(3.0),
                 framework_text, font_size=9, color=DARK_GRAY)

    # Employee & Social Metrics
    _gs_section_title(slide, "Social Metrics", Inches(4.8))
    social_rows = [
        ["Full-Time Employees", f"{cd.employees:,}" if cd.employees else "N/A"],
        ["Revenue per Employee", f"{cs}{cd.revenue.iloc[0]/cd.employees:,.0f}" if cd.employees and cd.revenue is not None and len(cd.revenue) > 0 and cd.employees > 0 else "N/A"],
    ]
    _gs_table(slide, ["Metric", "Value"], social_rows,
              Inches(6.8), Inches(5.15), Inches(6), Inches(0.8),
              col_widths=[Inches(3), Inches(3)])


def _company_slide_lbo(prs, cd: CompanyData, slide_num=None):
    """LBO Returns Analysis slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    cs = cd.currency_symbol

    _gs_header(slide, "LBO Returns Analysis", f"{cd.name} ({cd.ticker})")
    _gs_footer(slide, cd.ticker, slide_num=slide_num)

    # LBO Assumptions
    _gs_section_title(slide, "LBO Assumptions", Inches(1.0))

    ev = cd.enterprise_value or 0
    ebitda_val = cd.ebitda.iloc[0] if cd.ebitda is not None and len(cd.ebitda) > 0 else 0
    entry_multiple = (ev / ebitda_val) if ebitda_val and ebitda_val > 0 else 0
    leverage = 5.0  # Assumed 5x leverage
    equity_pct = max(0, ev - leverage * ebitda_val) / ev * 100 if ev > 0 and ebitda_val > 0 else 40

    assumptions = [
        ["Enterprise Value", format_number(ev, currency_symbol=cs)],
        ["Entry EV/EBITDA", f"{entry_multiple:.1f}x" if entry_multiple else "N/A"],
        ["LTM EBITDA", format_number(ebitda_val, currency_symbol=cs)],
        ["Assumed Leverage", f"{leverage:.1f}x EBITDA"],
        ["Equity Contribution", f"~{equity_pct:.0f}% of EV"],
    ]
    _gs_table(slide, ["Parameter", "Value"], assumptions,
              Inches(0.5), Inches(1.35), Inches(5.8), Inches(1.8),
              col_widths=[Inches(2.8), Inches(3.0)])

    # Returns Sensitivity
    _gs_section_title(slide, "Returns Sensitivity (5-Year Hold)", Inches(3.5))

    exit_multiples = [entry_multiple * 0.8, entry_multiple, entry_multiple * 1.2] if entry_multiple > 0 else [8, 10, 12]
    ebitda_growth_rates = [0.05, 0.08, 0.12]

    returns_rows = []
    for growth in ebitda_growth_rates:
        row = [f"{growth*100:.0f}% EBITDA Growth"]
        for exit_m in exit_multiples:
            exit_ebitda = ebitda_val * (1 + growth) ** 5 if ebitda_val > 0 else 0
            exit_ev = exit_ebitda * exit_m
            equity_in = ev * (equity_pct / 100) if ev > 0 else 1
            debt_paydown = ebitda_val * 0.3 * 5 if ebitda_val > 0 else 0  # 30% FCF conversion
            equity_out = max(0, exit_ev - (leverage * ebitda_val - debt_paydown))
            moic = equity_out / equity_in if equity_in > 0 else 0
            irr = (moic ** (1/5) - 1) * 100 if moic > 0 else 0
            row.append(f"{irr:.1f}% / {moic:.1f}x")
        returns_rows.append(row)

    exit_labels = [f"{m:.1f}x" for m in exit_multiples]
    _gs_table(slide, ["Scenario"] + [f"Exit {l}" for l in exit_labels], returns_rows,
              Inches(0.5), Inches(3.85), Inches(12.333), Inches(1.2),
              col_widths=[Inches(3)] + [Inches(3.111)] * 3)

    # Key LBO Considerations (right)
    _gs_section_title(slide, "Key LBO Considerations", Inches(1.0))

    fcf_yield = None
    if cd.free_cashflow_series is not None and len(cd.free_cashflow_series) > 0 and ev > 0:
        fcf_yield = cd.free_cashflow_series.iloc[0] / ev * 100

    considerations = [
        ["FCF Yield", f"{fcf_yield:.1f}%" if fcf_yield else "N/A"],
        ["Interest Coverage", f"{cd.interest_coverage:.1f}x" if cd.interest_coverage else "N/A"],
        ["Debt / EBITDA (Current)", f"{cd.net_debt_to_ebitda:.1f}x" if cd.net_debt_to_ebitda else "N/A"],
        ["Capital Intensity", "High" if (cd.capital_expenditure is not None and len(cd.capital_expenditure) > 0 and ebitda_val > 0 and abs(cd.capital_expenditure.iloc[0]) / ebitda_val > 0.3) else "Moderate"],
    ]
    _gs_table(slide, ["Factor", "Value"], considerations,
              Inches(6.8), Inches(1.35), Inches(6), Inches(1.5),
              col_widths=[Inches(3), Inches(3)])

    # Verdict
    _gs_section_title(slide, "LBO Viability", Inches(5.3))
    viable = "ATTRACTIVE" if (fcf_yield and fcf_yield > 5 and entry_multiple and entry_multiple < 12) else "MODERATE" if entry_multiple and entry_multiple < 15 else "CHALLENGING"
    v_color = GREEN if viable == "ATTRACTIVE" else (RGBColor(0xF5, 0xA6, 0x23) if viable == "MODERATE" else RED)
    _add_textbox(slide, Inches(0.5), Inches(5.65), Inches(5), Inches(0.5),
                 f"LBO Candidacy: {viable}", font_size=14, bold=True, color=v_color)


def _company_slide_mgmt(prs, cd: CompanyData, slide_num=None):
    """Management Effectiveness slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    cs = cd.currency_symbol

    _gs_header(slide, "Management Effectiveness", f"{cd.name} ({cd.ticker})")
    _gs_footer(slide, cd.ticker, slide_num=slide_num)

    # Returns on Capital
    _gs_section_title(slide, "Returns on Capital", Inches(1.0))

    roe = cd.return_on_equity
    roa = getattr(cd, 'return_on_assets', None)
    roic = getattr(cd, 'return_on_invested_capital', None)

    returns_rows = [
        ["Return on Equity (ROE)", f"{roe*100:.1f}%" if roe else "N/A"],
        ["Return on Assets (ROA)", f"{roa*100:.1f}%" if roa else "N/A"],
        ["Return on Invested Capital (ROIC)", f"{roic*100:.1f}%" if roic else "N/A"],
    ]
    _gs_table(slide, ["Metric", "Value"], returns_rows,
              Inches(0.5), Inches(1.35), Inches(5.8), Inches(1.1),
              col_widths=[Inches(3.3), Inches(2.5)])

    # Efficiency Metrics
    _gs_section_title(slide, "Operational Efficiency", Inches(2.7))

    asset_turnover = getattr(cd, 'asset_turnover', None)
    inventory_turnover = getattr(cd, 'inventory_turnover', None)
    receivables_turnover = getattr(cd, 'receivables_turnover', None)

    eff_rows = [
        ["Asset Turnover", f"{asset_turnover:.2f}x" if asset_turnover else "N/A"],
        ["Inventory Turnover", f"{inventory_turnover:.1f}x" if inventory_turnover else "N/A"],
        ["Receivables Turnover", f"{receivables_turnover:.1f}x" if receivables_turnover else "N/A"],
        ["Gross Margin", f"{cd.gross_margins*100:.1f}%" if cd.gross_margins else "N/A"],
        ["Operating Margin", f"{cd.operating_margins*100:.1f}%" if cd.operating_margins else "N/A"],
    ]
    _gs_table(slide, ["Metric", "Value"], eff_rows,
              Inches(0.5), Inches(3.05), Inches(5.8), Inches(1.8),
              col_widths=[Inches(3.3), Inches(2.5)])

    # DuPont Decomposition (right side)
    _gs_section_title(slide, "DuPont Analysis (ROE Decomposition)", Inches(1.0))

    net_margin = cd.profit_margins or 0
    at = asset_turnover or 0
    # Equity multiplier approximation
    total_assets_val = cd.total_assets.iloc[0] if cd.total_assets is not None and len(cd.total_assets) > 0 else 0
    total_equity_val = cd.total_equity.iloc[0] if cd.total_equity is not None and len(cd.total_equity) > 0 else 0
    equity_mult = total_assets_val / total_equity_val if total_equity_val and total_equity_val > 0 else 0

    dupont_rows = [
        ["Net Profit Margin", f"{net_margin*100:.1f}%" if net_margin else "N/A"],
        ["× Asset Turnover", f"{at:.2f}x" if at else "N/A"],
        ["× Equity Multiplier", f"{equity_mult:.2f}x" if equity_mult else "N/A"],
        ["= ROE", f"{net_margin * at * equity_mult * 100:.1f}%" if all([net_margin, at, equity_mult]) else "N/A"],
    ]
    _gs_table(slide, ["Component", "Value"], dupont_rows,
              Inches(6.8), Inches(1.35), Inches(6), Inches(1.5),
              col_widths=[Inches(3), Inches(3)])

    # Capital Allocation
    _gs_section_title(slide, "Capital Allocation", Inches(3.0))

    div_yield = cd.dividend_yield
    payout = getattr(cd, 'payout_ratio', None)
    buyback = getattr(cd, 'share_buyback_yield', None)

    alloc_rows = [
        ["Dividend Yield", f"{div_yield*100:.2f}%" if div_yield else "N/A"],
        ["Payout Ratio", f"{payout*100:.1f}%" if payout else "N/A"],
        ["CapEx / Revenue", f"{abs(cd.capital_expenditure.iloc[0]) / cd.revenue.iloc[0] * 100:.1f}%" if cd.capital_expenditure is not None and len(cd.capital_expenditure) > 0 and cd.revenue is not None and len(cd.revenue) > 0 and cd.revenue.iloc[0] > 0 else "N/A"],
    ]
    _gs_table(slide, ["Metric", "Value"], alloc_rows,
              Inches(6.8), Inches(3.35), Inches(6), Inches(1.1),
              col_widths=[Inches(3), Inches(3)])

    # Piotroski Score
    _gs_section_title(slide, "Financial Health Score", Inches(5.0))
    try:
        piotroski = calculate_piotroski_score(cd)
        if piotroski:
            score = piotroski['score']
            score_color = GREEN if score >= 7 else (RED if score <= 3 else DARK_GRAY)
            label = "Strong" if score >= 7 else ("Weak" if score <= 3 else "Neutral")
            _add_textbox(slide, Inches(0.5), Inches(5.35), Inches(5), Inches(0.5),
                         f"Piotroski F-Score: {score} / 9  ({label})", font_size=16, bold=True, color=score_color)
    except Exception:
        _add_textbox(slide, Inches(0.5), Inches(5.35), Inches(5), Inches(0.5),
                     "Piotroski F-Score: N/A", font_size=12, color=TEXT_DIM)


def generate_presentation(cd: CompanyData, template_path: str = "assets/template.pptx", confidential: bool = False) -> io.BytesIO:
    """Build the 7-slide dark-theme company profile with purple accents."""
    prs = Presentation(template_path)
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    _company_slide_1(prs, cd)           # Executive Summary
    _company_slide_2(prs, cd)           # Financial Overview
    _company_slide_3(prs, cd)           # Peer Comparison
    _company_slide_valuation(prs, cd)   # Valuation Summary
    _company_slide_esg(prs, cd)         # ESG Summary
    _company_slide_lbo(prs, cd)         # LBO Returns
    _company_slide_mgmt(prs, cd)        # Management Effectiveness

    # Add slide numbers to all slides
    for i, slide in enumerate(prs.slides, start=1):
        _add_textbox(slide, Inches(6), SLIDE_H - Inches(0.4), Inches(1.333), Inches(0.3),
                     str(i), font_size=7, color=TEXT_DIM, align=PP_ALIGN.CENTER)

    # Add confidential watermark if requested
    if confidential:
        for slide in prs.slides:
            _add_confidential_watermark(slide)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf


# ══════════════════════════════════════════════════════════════
# MERGER DEAL BOOK - 3 SLIDES
# ══════════════════════════════════════════════════════════════

def _deal_slide_1(prs, acq, tgt, pf, assumptions):
    """Slide 1: Transaction Overview."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    cs = acq.currency_symbol

    _gs_header(slide, "Transaction Overview", f"{acq.name} to acquire {tgt.name}")
    _gs_footer(slide, f"{acq.ticker} + {tgt.ticker}")

    # Key Transaction Terms (left)
    _gs_section_title(slide, "Key Transaction Terms", Inches(1.0))

    terms = [
        ["Purchase Price", format_number(pf.purchase_price, currency_symbol=cs)],
        ["Offer Premium", f"{assumptions.offer_premium_pct:.0f}%"],
        ["Offer Price / Share", f"{cs}{pf.offer_price_per_share:.2f}" if pf.offer_price_per_share else "—"],
        ["Implied EV/EBITDA", f"{pf.implied_ev_ebitda:.1f}x" if pf.implied_ev_ebitda else "—"],
        ["Implied P/E", f"{pf.implied_pe:.1f}x" if pf.implied_pe else "—"],
        ["Transaction Fees", format_number(pf.transaction_fees, currency_symbol=cs)],
    ]
    _gs_table(slide, ["Metric", "Value"], terms,
              Inches(0.5), Inches(1.35), Inches(5.5), Inches(2.2),
              col_widths=[Inches(2.5), Inches(3)])

    # Consideration Mix (left bottom)
    _gs_section_title(slide, "Consideration Structure", Inches(3.8))

    consideration = [
        ["Cash Consideration", format_number(pf.cash_consideration, currency_symbol=cs), f"{assumptions.pct_cash:.0f}%"],
        ["Stock Consideration", format_number(pf.stock_consideration, currency_symbol=cs), f"{assumptions.pct_stock:.0f}%"],
        ["New Shares Issued", f"{pf.new_shares_issued/1e6:.1f}M", "—"],
        ["Total Consideration", format_number(pf.purchase_price, currency_symbol=cs), "100%"],
    ]
    _gs_table(slide, ["Component", "Amount", "% of Total"], consideration,
              Inches(0.5), Inches(4.15), Inches(5.5), Inches(1.5),
              col_widths=[Inches(2.2), Inches(2.0), Inches(1.3)])

    # Company Comparison (right)
    _gs_section_title(slide, "Company Comparison", Inches(1.0))

    comparison = [
        ["Market Cap", format_number(acq.market_cap, currency_symbol=cs), format_number(tgt.market_cap, currency_symbol=cs)],
        ["Enterprise Value", format_number(acq.enterprise_value, currency_symbol=cs), format_number(tgt.enterprise_value, currency_symbol=cs)],
        ["Revenue (LTM)", format_number(pf.acq_revenue, currency_symbol=cs), format_number(pf.tgt_revenue, currency_symbol=cs)],
        ["EBITDA (LTM)", format_number(pf.acq_ebitda, currency_symbol=cs), format_number(pf.tgt_ebitda, currency_symbol=cs)],
        ["Net Income", format_number(pf.acq_net_income, currency_symbol=cs), format_number(pf.tgt_net_income, currency_symbol=cs)],
        ["EV/EBITDA", f"{acq.ev_to_ebitda:.1f}x" if acq.ev_to_ebitda else "—", f"{tgt.ev_to_ebitda:.1f}x" if tgt.ev_to_ebitda else "—"],
    ]
    _gs_table(slide, ["Metric", acq.ticker, tgt.ticker], comparison,
              Inches(6.8), Inches(1.35), Inches(6), Inches(2.2),
              col_widths=[Inches(2.2), Inches(1.9), Inches(1.9)])

    # Deal Rationale (right bottom)
    _gs_section_title(slide, "Strategic Rationale", Inches(3.8))

    rationale = [
        ["Cost Synergies", format_number(pf.cost_synergies, currency_symbol=cs)],
        ["Revenue Synergies", format_number(pf.revenue_synergies, currency_symbol=cs)],
        ["Total Synergies", format_number(pf.total_synergies, currency_symbol=cs)],
        ["Synergies % of Target Rev", f"{pf.total_synergies/pf.tgt_revenue*100:.1f}%" if pf.tgt_revenue else "—"],
    ]
    _gs_table(slide, ["Metric", "Value"], rationale,
              Inches(6.8), Inches(4.15), Inches(6), Inches(1.5),
              col_widths=[Inches(3.5), Inches(2.5)])


def _deal_slide_2(prs, acq, tgt, pf, assumptions):
    """Slide 2: Financial Impact - Pro Forma & Accretion/Dilution."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    cs = acq.currency_symbol
    tax_r = assumptions.tax_rate / 100

    _gs_header(slide, "Financial Impact Analysis", f"{acq.ticker} + {tgt.ticker}")
    _gs_footer(slide, f"{acq.ticker} + {tgt.ticker}")

    # Pro Forma P&L (left)
    _gs_section_title(slide, "Pro Forma Income Statement", Inches(1.0))

    ats = pf.total_synergies * (1 - tax_r)
    ati = pf.incremental_interest * (1 - tax_r)

    pf_rows = [
        ["Revenue", format_number(pf.acq_revenue, currency_symbol=cs), format_number(pf.tgt_revenue, currency_symbol=cs), format_number(pf.revenue_synergies, currency_symbol=cs), format_number(pf.pf_revenue, currency_symbol=cs)],
        ["EBITDA", format_number(pf.acq_ebitda, currency_symbol=cs), format_number(pf.tgt_ebitda, currency_symbol=cs), format_number(pf.total_synergies, currency_symbol=cs), format_number(pf.pf_ebitda, currency_symbol=cs)],
        ["Net Income", format_number(pf.acq_net_income, currency_symbol=cs), format_number(pf.tgt_net_income, currency_symbol=cs), format_number(ats - ati, currency_symbol=cs), format_number(pf.pf_net_income, currency_symbol=cs)],
        ["Shares (M)", f"{pf.acq_shares/1e6:.0f}" if pf.acq_shares else "—", "—", f"+{pf.new_shares_issued/1e6:.0f}" if pf.new_shares_issued else "—", f"{pf.pf_shares_outstanding/1e6:.0f}" if pf.pf_shares_outstanding else "—"],
        ["EPS", f"{cs}{pf.acq_eps:.2f}" if pf.acq_eps else "—", "—", "—", f"{cs}{pf.pf_eps:.2f}" if pf.pf_eps else "—"],
    ]
    _gs_table(slide, ["", acq.ticker, tgt.ticker, "Adj.", "Pro Forma"], pf_rows,
              Inches(0.5), Inches(1.35), Inches(7.5), Inches(1.9),
              col_widths=[Inches(1.5), Inches(1.5), Inches(1.5), Inches(1.5), Inches(1.5)])

    # Accretion/Dilution Summary (left bottom)
    _gs_section_title(slide, "Accretion / Dilution", Inches(3.5))

    acc_color = GREEN if pf.is_accretive else RED
    acc_word = "ACCRETIVE" if pf.is_accretive else "DILUTIVE"

    # Big impact number
    _add_textbox(slide, Inches(0.5), Inches(3.85), Inches(3), Inches(0.6),
                 f"{pf.accretion_dilution_pct:+.1f}%", font_size=28, bold=True, color=acc_color)
    _add_textbox(slide, Inches(0.5), Inches(4.4), Inches(3), Inches(0.3),
                 acc_word, font_size=12, bold=True, color=acc_color)

    # EPS bridge table
    eps_rows = [
        ["Standalone EPS", f"{cs}{pf.acq_eps:.2f}" if pf.acq_eps else "—"],
        ["Pro Forma EPS", f"{cs}{pf.pf_eps:.2f}" if pf.pf_eps else "—"],
        ["EPS Change", f"{cs}{pf.pf_eps - pf.acq_eps:.2f}" if pf.acq_eps and pf.pf_eps else "—"],
    ]
    _gs_table(slide, ["Metric", "Value"], eps_rows,
              Inches(3.8), Inches(3.85), Inches(3.7), Inches(1.1),
              col_widths=[Inches(2), Inches(1.7)])

    # Sources & Uses (right)
    _gs_section_title(slide, "Sources of Funds", Inches(1.0))

    sources_rows = [[k, format_number(v, currency_symbol=cs)] for k, v in pf.sources.items()]
    _gs_table(slide, ["Source", "Amount"], sources_rows,
              Inches(8.5), Inches(1.35), Inches(4.333), Inches(1.5),
              col_widths=[Inches(2.5), Inches(1.833)])

    _gs_section_title(slide, "Uses of Funds", Inches(3.0))

    uses_rows = [[k, format_number(v, currency_symbol=cs)] for k, v in pf.uses.items()]
    _gs_table(slide, ["Use", "Amount"], uses_rows,
              Inches(8.5), Inches(3.35), Inches(4.333), Inches(1.5),
              col_widths=[Inches(2.5), Inches(1.833)])

    # Credit Metrics (bottom right)
    _gs_section_title(slide, "Pro Forma Credit Profile", Inches(5.0))

    credit_rows = [
        ["PF Debt / EBITDA", f"{pf.pf_leverage_ratio:.1f}x" if pf.pf_leverage_ratio else "—"],
        ["PF Interest Coverage", f"{pf.pf_interest_coverage:.1f}x" if pf.pf_interest_coverage else "—"],
        ["PF Total Debt", format_number(pf.pf_total_debt, currency_symbol=cs)],
    ]
    _gs_table(slide, ["Metric", "Value"], credit_rows,
              Inches(8.5), Inches(5.35), Inches(4.333), Inches(1.1),
              col_widths=[Inches(2.5), Inches(1.833)])


def _deal_slide_3(prs, acq, tgt, pf, assumptions, football_field):
    """Slide 3: Valuation Analysis - Football Field."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    cs = acq.currency_symbol

    _gs_header(slide, "Valuation Analysis", f"{tgt.name} ({tgt.ticker})")
    _gs_footer(slide, f"{acq.ticker} + {tgt.ticker}")

    # Football Field Table
    _gs_section_title(slide, "Valuation Summary (Football Field)", Inches(1.0))

    offer_price = football_field.get("_offer_price", 0) if football_field else 0
    methods = {k: v for k, v in football_field.items() if not k.startswith("_")} if football_field else {}

    def _fmt_val(v):
        if v is None:
            return "—"
        if abs(v) >= 1e9:
            return f"{cs}{v/1e9:.1f}B"
        elif abs(v) >= 1e6:
            return f"{cs}{v/1e6:.0f}M"
        else:
            return f"{cs}{v:,.0f}"

    ff_rows = []
    for method, vals in methods.items():
        low = vals.get("low", 0)
        high = vals.get("high", 0)
        mid = (low + high) / 2 if low and high else 0
        ff_rows.append([method, _fmt_val(low), _fmt_val(mid), _fmt_val(high)])

    if ff_rows:
        _gs_table(slide, ["Methodology", "Low", "Midpoint", "High"], ff_rows,
                  Inches(0.5), Inches(1.35), Inches(7), Inches(2.2),
                  col_widths=[Inches(2.5), Inches(1.5), Inches(1.5), Inches(1.5)])
    else:
        _add_textbox(slide, Inches(0.5), Inches(1.35), Inches(7), Inches(0.5),
                     "Insufficient data for valuation analysis", font_size=10, color=DARK_GRAY)

    # Offer Price Comparison
    _gs_section_title(slide, "Offer Analysis", Inches(3.8))

    offer_rows = [
        ["Offer Price / Share", f"{cs}{pf.offer_price_per_share:.2f}" if pf.offer_price_per_share else "—"],
        ["Current Price", f"{cs}{tgt.current_price:.2f}" if tgt.current_price else "—"],
        ["Offer Premium", f"{assumptions.offer_premium_pct:.0f}%"],
        ["52-Week High", f"{cs}{tgt.fifty_two_week_high:.2f}" if tgt.fifty_two_week_high else "—"],
        ["52-Week Low", f"{cs}{tgt.fifty_two_week_low:.2f}" if tgt.fifty_two_week_low else "—"],
    ]
    _gs_table(slide, ["Metric", "Value"], offer_rows,
              Inches(0.5), Inches(4.15), Inches(5), Inches(1.8),
              col_widths=[Inches(2.5), Inches(2.5)])

    # Implied Multiples (right side)
    _gs_section_title(slide, "Implied Transaction Multiples", Inches(1.0))

    multiples = [
        ["Implied EV / Revenue", f"{pf.purchase_price / pf.tgt_revenue:.1f}x" if pf.tgt_revenue else "—"],
        ["Implied EV / EBITDA", f"{pf.implied_ev_ebitda:.1f}x" if pf.implied_ev_ebitda else "—"],
        ["Implied P / E", f"{pf.implied_pe:.1f}x" if pf.implied_pe else "—"],
    ]
    _gs_table(slide, ["Multiple", "Value"], multiples,
              Inches(8), Inches(1.35), Inches(4.833), Inches(1.2),
              col_widths=[Inches(2.8), Inches(2.033)])

    # Transaction Assumptions
    _gs_section_title(slide, "Transaction Assumptions", Inches(2.8))

    assump_rows = [
        ["Cost of Debt", f"{assumptions.cost_of_debt:.1f}%"],
        ["Tax Rate", f"{assumptions.tax_rate:.1f}%"],
        ["Cost Synergies (% of Target SG&A)", f"{assumptions.cost_synergies_pct:.1f}%"],
        ["Revenue Synergies (% of Target Rev)", f"{assumptions.revenue_synergies_pct:.1f}%"],
        ["Transaction Fees (% of Deal)", f"{assumptions.transaction_fees_pct:.1f}%"],
    ]
    _gs_table(slide, ["Assumption", "Value"], assump_rows,
              Inches(8), Inches(3.15), Inches(4.833), Inches(1.9),
              col_widths=[Inches(3.3), Inches(1.533)])


def generate_deal_book(acq_cd, tgt_cd, pro_forma, merger_insights, assumptions,
                       template_path: str = "assets/template.pptx", confidential: bool = False) -> io.BytesIO:
    """Build the 3-slide dark-theme deal book with purple accents."""
    prs = Presentation(template_path)
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    _deal_slide_1(prs, acq_cd, tgt_cd, pro_forma, assumptions)  # Transaction Overview
    _deal_slide_2(prs, acq_cd, tgt_cd, pro_forma, assumptions)  # Financial Impact
    _deal_slide_3(prs, acq_cd, tgt_cd, pro_forma, assumptions, pro_forma.football_field)  # Valuation

    # Add slide numbers
    for i, slide in enumerate(prs.slides, start=1):
        _add_textbox(slide, Inches(6), SLIDE_H - Inches(0.4), Inches(1.333), Inches(0.3),
                     str(i), font_size=7, color=TEXT_DIM, align=PP_ALIGN.CENTER)

    if confidential:
        for slide in prs.slides:
            _add_confidential_watermark(slide)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf
