"""
Generate template.pptx — a branded 3-layout landscape template.

Run once:  python create_template.py
Produces:  assets/template.pptx

Each slide layout has indexed placeholders that the generation engine
references by idx.  Run `python template_inspector.py` after creation
to verify the mapping.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from lxml import etree
import os

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

NAVY = RGBColor(0x0B, 0x1D, 0x3A)
GOLD = RGBColor(0xD4, 0xA5, 0x37)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)


def _add_ph(sp_tree, idx, left, top, width, height, name, ph_type=None):
    """Add a placeholder <p:sp> element to a slide layout's shape tree."""
    sp = etree.SubElement(sp_tree, qn("p:sp"))

    # -- nvSpPr --
    nvSpPr = etree.SubElement(sp, qn("p:nvSpPr"))
    cNvPr = etree.SubElement(nvSpPr, qn("p:cNvPr"))
    cNvPr.set("id", str(idx + 200))
    cNvPr.set("name", name)

    cNvSpPr = etree.SubElement(nvSpPr, qn("p:cNvSpPr"))
    locks = etree.SubElement(cNvSpPr, qn("a:spLocks"))
    locks.set("noGrp", "1")

    nvPr = etree.SubElement(nvSpPr, qn("p:nvPr"))
    ph = etree.SubElement(nvPr, qn("p:ph"))
    ph.set("idx", str(idx))
    if ph_type:
        ph.set("type", ph_type)

    # -- spPr --
    spPr = etree.SubElement(sp, qn("p:spPr"))
    xfrm = etree.SubElement(spPr, qn("a:xfrm"))
    off = etree.SubElement(xfrm, qn("a:off"))
    off.set("x", str(int(Emu(left) if isinstance(left, int) else left)))
    off.set("y", str(int(Emu(top) if isinstance(top, int) else top)))
    ext = etree.SubElement(xfrm, qn("a:ext"))
    ext.set("cx", str(int(Emu(width) if isinstance(width, int) else width)))
    ext.set("cy", str(int(Emu(height) if isinstance(height, int) else height)))

    # -- txBody --
    txBody = etree.SubElement(sp, qn("p:txBody"))
    bodyPr = etree.SubElement(txBody, qn("a:bodyPr"))
    bodyPr.set("wrap", "square")
    etree.SubElement(txBody, qn("a:lstStyle"))
    etree.SubElement(txBody, qn("a:p"))

    return sp


def _get_sp_tree(layout):
    return layout._element.find(qn("p:cSld")).find(qn("p:spTree"))


def build():
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    # We modify the built-in Blank layout (index 6) and two others
    # to hold our indexed placeholders.

    # ── LAYOUT 0 — Executive Summary ────────────────────────
    # Use layout index 6 (Blank)
    lay0 = prs.slide_layouts[6]
    lay0.name = "Executive Summary"
    tree0 = _get_sp_tree(lay0)

    # idx 10 — Company Name (top-left title area)
    _add_ph(tree0, 10, Inches(0.5), Inches(0.3), Inches(6), Inches(0.7), "Company Name")
    # idx 11 — Ticker / Price line
    _add_ph(tree0, 11, Inches(0.5), Inches(1.0), Inches(6), Inches(0.5), "Ticker Price")
    # idx 12 — 3-bullet summary
    _add_ph(tree0, 12, Inches(0.5), Inches(1.7), Inches(5.5), Inches(4.5), "Summary Bullets")
    # idx 13 — Price chart image placeholder
    _add_ph(tree0, 13, Inches(6.8), Inches(0.3), Inches(6), Inches(4.0), "Price Chart")
    # idx 14 — Key metrics box
    _add_ph(tree0, 14, Inches(6.8), Inches(4.5), Inches(6), Inches(2.5), "Key Metrics")

    # ── LAYOUT 1 — Financials & Deal Score ──────────────────
    lay1 = prs.slide_layouts[5]  # Title Only
    lay1.name = "Financials"
    tree1 = _get_sp_tree(lay1)

    # idx 20 — Slide title
    _add_ph(tree1, 20, Inches(0.5), Inches(0.3), Inches(12), Inches(0.7), "Fin Title")
    # idx 21 — Financial table area
    _add_ph(tree1, 21, Inches(0.5), Inches(1.3), Inches(6), Inches(5.5), "Fin Table")
    # idx 22 — EBITDA margin chart area
    _add_ph(tree1, 22, Inches(7.0), Inches(1.3), Inches(5.8), Inches(3.5), "EBITDA Chart")
    # idx 23 — Deal score area
    _add_ph(tree1, 23, Inches(7.0), Inches(5.0), Inches(5.8), Inches(2.0), "Deal Score")

    # ── LAYOUT 2 — Strategy & M&A ──────────────────────────
    lay2 = prs.slide_layouts[4]  # Two Content
    lay2.name = "Strategy"
    tree2 = _get_sp_tree(lay2)

    # idx 30 — Slide title
    _add_ph(tree2, 30, Inches(0.5), Inches(0.3), Inches(12), Inches(0.7), "Strat Title")
    # idx 31 — Management names
    _add_ph(tree2, 31, Inches(0.5), Inches(1.3), Inches(5.5), Inches(2.5), "Management")
    # idx 32 — Product segments / pie chart area
    _add_ph(tree2, 32, Inches(7.0), Inches(1.3), Inches(5.8), Inches(3.0), "Segments Pie")
    # idx 33 — News headlines
    _add_ph(tree2, 33, Inches(0.5), Inches(4.2), Inches(12), Inches(3.0), "News")

    os.makedirs("assets", exist_ok=True)
    prs.save("assets/template.pptx")
    print("assets/template.pptx created successfully.")


if __name__ == "__main__":
    build()
